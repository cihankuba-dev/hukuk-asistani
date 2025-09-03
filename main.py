import os
import io
import tempfile
import pickle
import numpy as np
import faiss

from fastapi import FastAPI, UploadFile, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from openai import OpenAI

from PyPDF2 import PdfReader
import httpx
from bs4 import BeautifulSoup
import docx
import openpyxl
from pptx import Presentation

# Google Drive (opsiyonel – credentials.json yoksa ingest_drive çalışmaz)
from google.oauth2 import service_account
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

# ==============
# OpenAI Client
# ==============
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# ===================
# FAISS (Arşiv/Hafıza)
# ===================
INDEX_FILE = "faiss_index.pkl"
EMBED_DIM = 1536  # text-embedding-3-small

if os.path.exists(INDEX_FILE):
    with open(INDEX_FILE, "rb") as f:
        index, metadata = pickle.load(f)
else:
    index = faiss.IndexFlatL2(EMBED_DIM)
    metadata = []  # [{"text": "...", "file": "xxx.pdf"}]

# =========
# FastAPI
# =========
app = FastAPI(title="⚖️ Hukuk Asistanı")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/")
async def root():
    return JSONResponse(content={"status": "ok", "message": "⚖️ Hukuk Asistanı aktif ve çalışıyor!"})

@app.get("/health")
async def health():
    return JSONResponse(content={"ok": True})

# =========================
# Yardımcı: Embedding & Chat
# =========================
def embed_text(text: str):
    resp = client.embeddings.create(
        model="text-embedding-3-small",
        input=text
    )
    return np.array(resp.data[0].embedding, dtype="float32")

def chat(messages, max_tokens=1500):
    resp = client.chat.completions.create(
        model="gpt-5",
        messages=messages,
        max_completion_tokens=max_tokens
        # ❌ temperature parametresini kaldırdık
    )
    return resp.choices[0].message.content

# =========================
# Mevzuat ve İçtihat Modülleri
# =========================
async def fetch_mevzuat(query: str):
    url = f"https://www.mevzuat.gov.tr/arama?aranan={query}"
    try:
        async with httpx.AsyncClient() as ac:
            r = await ac.get(url, timeout=60)
            r.encoding = "utf-8"
            soup = BeautifulSoup(r.text, "html.parser")
            results = [a.text.strip() for a in soup.select("a") if a.text.strip()]
        return results[:5] if results else ["Sonuç bulunamadı."]
    except Exception as e:
        return [f"Mevzuat bilgisi alınamadı: {e}"]

async def search_ictihat(keyword: str, limit: int = 3) -> list[str]:
    url = f"https://karararama.yargitay.gov.tr/Yargitay-Karar-Forumu?q={keyword}"
    try:
        async with httpx.AsyncClient() as ac:
            resp = await ac.get(url, timeout=100)
            resp.encoding = "utf-8"
            if resp.status_code != 200:
                return [f"Emsal karar bulunamadı ({keyword})."]
            soup = BeautifulSoup(resp.text, "html.parser")
            results = [div.get_text(strip=True) for div in soup.select("div.kararOzet")[:limit]]
            return results if results else [f"Emsal karar bulunamadı ({keyword})."]
    except Exception as e:
        return [f"İçtihat bilgisi alınamadı: {e}"]

# =========================
# Dosya Metin Çıkarma
# =========================
def extract_text_from_path(path, filename):
    ext = filename.split(".")[-1].lower()
    text = ""
    try:
        if ext == "pdf":
            reader = PdfReader(path)
            text = "\n".join([(page.extract_text() or "") for page in reader.pages])
        elif ext == "docx":
            doc = docx.Document(path)
            text = "\n".join([p.text for p in doc.paragraphs])
        elif ext == "xlsx":
            wb = openpyxl.load_workbook(path, data_only=True)
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    row_txt = " ".join([str(cell) for cell in row if cell is not None]).strip()
                    if row_txt:
                        text += row_txt + "\n"
        elif ext == "pptx":
            prs = Presentation(path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
        elif ext in ["txt", "rtf", "md", "udf"]:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
    except Exception as e:
        print(f"[HATA] {filename}: {e}")
    return (text or "").strip()

# =========================
# Google Drive servis bağlantısı
# =========================
def get_drive_service():
    creds = service_account.Credentials.from_service_account_file(
        "/etc/secrets/credentials.json",   # ✅ Render'daki doğru path
        scopes=["https://www.googleapis.com/auth/drive.readonly"],
    )
    return build("drive", "v3", credentials=creds)

# =========================
# ENDPOINTS
# =========================

@app.get("/test_key")
async def test_key():
    try:
        out = chat(
            messages=[{"role": "user", "content": "Merhaba, sadece test ediyorum."}],
            max_tokens=50
        )
        return JSONResponse(content={"status": "ok", "response": out})
    except Exception as e:
        return JSONResponse(content={"status": "error", "message": str(e)})

@app.post("/ingest")
async def ingest(file: UploadFile):
    """Tek dosya yükleyip arşive (FAISS) ekler."""
    global index, metadata
    try:
        ext = file.filename.split(".")[-1].lower()
        content = await file.read()
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        text = extract_text_from_path(tmp_path, file.filename)
        os.remove(tmp_path)

        if not text:
            return JSONResponse(content={"status": "error", "message": "Metin çıkarılamadı."})

        # 1000 karakterlik chunk'lara böl
        chunks = [text[i:i+1000] for i in range(0, len(text), 1000)]
        added = 0
        for chunk in chunks:
            vec = embed_text(chunk)
            index.add(np.array([vec], dtype="float32"))
            metadata.append({"text": chunk, "file": file.filename})
            added += 1

        with open(INDEX_FILE, "wb") as f:
            pickle.dump((index, metadata), f)

        return JSONResponse(content={"status": "ok", "file": file.filename, "chunks_added": added, "chunks_total": len(metadata)})
    except Exception as e:
        return JSONResponse(content={"status": "error", "message": str(e)})

@app.post("/ingest_drive")
async def ingest_drive(folder_id: str = Form(...)):
    """Google Drive klasöründen desteklenen dosyaları indirip arşive ekler."""
    global index, metadata
    try:
        service = get_drive_service()
        results = service.files().list(
            q=f"'{folder_id}' in parents and trashed=false",
            fields="files(id, name, mimeType)"
        ).execute()

        files = results.get("files", [])
        count_files, count_chunks = 0, 0

        for fobj in files:
            fname = fobj["name"]
            ext = fname.split(".")[-1].lower()
            if ext not in ["pdf", "docx", "xlsx", "pptx", "txt", "rtf", "md"]:
                continue

            request = service.files().get_media(fileId=fobj["id"])
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            fh.seek(0)

            temp_path = os.path.join(tempfile.gettempdir(), fname)
            with open(temp_path, "wb") as wf:
                wf.write(fh.read())

            text = extract_text_from_path(temp_path, fname)
            os.remove(temp_path)

            if not text:
                continue

            chunks = [text[i:i+1000] for i in range(0, len(text), 1000)]
            for ch in chunks:
                vec = embed_text(ch)
                index.add(np.array([vec], dtype="float32"))
                metadata.append({"text": ch, "file": fname})
                count_chunks += 1

            count_files += 1

        with open(INDEX_FILE, "wb") as f:
            pickle.dump((index, metadata), f)

        return JSONResponse(content={"status": "ok", "files_ingested": count_files, "chunks_total": len(metadata)})
    except Exception as e:
        return JSONResponse(content={"status": "error", "message": str(e)})

@app.post("/petition")
async def petition(prompt: str = Form(...)):
    """Serbest metinden dilekçe/talep metni üretir (üslup + imza içerir)."""
    try:
        messages = [
            {"role": "system", "content": "Sen deneyimli bir Türk hukuk asistanısın. Çıktının sonuna 'Av. Mehmet Cihan KUBA' imzasını ekle."},
            {"role": "user", "content": prompt}
        ]
        out = chat(messages, max_tokens=2000, temperature=0.2)
        return JSONResponse(content={"draft": out})
    except Exception as e:
        return JSONResponse(content={"error": str(e)})

@app.post("/summarize")
async def summarize(file: UploadFile):
    """Yüklenen dosyayı analiz edip özet çıkarır."""
    try:
        ext = file.filename.split(".")[-1].lower()
        content = await file.read()
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        text = extract_text_from_path(tmp_path, file.filename)
        os.remove(tmp_path)

        messages = [
            {"role": "system", "content": "Sen deneyimli bir hukuk asistanısın. Belgeleri analiz edip açık, maddeli özet çıkar."},
            {"role": "user", "content": f"Şu belgeyi özetle:\n\n{text[:6000]}"}  # güvenli sınır
        ]
        out = chat(messages, max_tokens=1000, temperature=0.2)
        return JSONResponse(content={"summary": out})
    except Exception as e:
        return JSONResponse(content={"error": str(e)})

@app.post("/draft_from_file")
async def draft_from_file(file: UploadFile, type: str = Form(...)):
    """
    Yüklenen dosyadan, arşiv (FAISS) + mevzuat + içtihat + üslup örnekleri kullanılarak
    istenen türde taslak üretir.
    """
    try:
        # 1) Dosya içeriği
        ext = file.filename.split(".")[-1].lower()
        content = await file.read()
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        text = extract_text_from_path(tmp_path, file.filename)
        os.remove(tmp_path)

        # 2) Arşiv: benzer içerik
        vec = embed_text(text)
        k = min(5, index.ntotal) if index.ntotal > 0 else 0
        context = ""
        if k > 0:
            D, I = index.search(np.array([vec], dtype="float32"), k=k)
            context = "\n".join([metadata[i]["text"] for i in I[0] if i < len(metadata)])

        # 3) Arşiv: üslup örnekleri (max 3 – .pdf/.docx’ten gelenler)
        style_examples = "\n\n".join(
            [m["text"] for m in metadata if m["file"].lower().endswith((".docx", ".pdf"))][:3]
        )

        # 4) Mevzuat & İçtihat
        mevzuat_bilgisi = await fetch_mevzuat(type)
        ictihatlar = await search_ictihat(type)

        # 5) Prompt
        messages = [
            {"role": "system", "content": "Sen deneyimli bir Türk hukuk asistanısın. Üslup resmi ve tutarlı olsun. Sonuna 'Av. Mehmet Cihan KUBA' imzası ekle."},
            {"role": "user", "content":
                f"Belge türü: {type}\n\n"
                f"Arşivden ilgili içerik (benzer pasajlar):\n{context}\n\n"
                f"Mevzuat bilgisi (ham arama sonuçları):\n{mevzuat_bilgisi}\n\n"
                f"İçtihat özetleri (ham arama sonuçları):\n{ictihatlar}\n\n"
                f"Üslup örnekleri (bunlara benzer yaz):\n{style_examples}\n\n"
                f"Dosya içeriği (ham metin):\n{text[:4000]}\n\n"
                f"Lütfen yukarıdaki tüm bilgilerden faydalanarak ayrıntılı bir {type} taslağı hazırla. "
                f"Kaynakları doğrudan link verme; metin içinde mevzuat/karar atıf biçiminde an ve sonunda kısa madde madde özet ver."}
        ]

        out = chat(messages, max_tokens=3000, temperature=0.2)
        return JSONResponse(content={"draft": out})
    except Exception as e:
        return JSONResponse(content={"error": str(e)})

@app.post("/law_search")
async def law_search(query: str = Form(...)):
    """Serbest sorgu: mevzuat + içtihat araması ve özet çıktısı."""
    try:
        mevzuat_bilgisi = await fetch_mevzuat(query)
        ictihatlar = await search_ictihat(query)

        messages = [
            {"role": "system", "content": "Sen deneyimli bir hukuk araştırma asistanısın. Güncel mevzuat ve içtihatlardan alıntılarla kısa özet ver; belirsizlik varsa açıkça belirt."},
            {"role": "user", "content": f"Sorgu: {query}\n\nMevzuat ham sonuçlar: {mevzuat_bilgisi}\n\nİçtihat ham sonuçlar: {ictihatlar}\n\nBunları derle ve kısa bir özet/yorum üret."}
        ]
        out = chat(messages, max_tokens=1500, temperature=0.2)
        return JSONResponse(content={"result": out, "mevzuat_raw": mevzuat_bilgisi, "ictihat_raw": ictihatlar})
    except Exception as e:
        return JSONResponse(content={"error": str(e)})
